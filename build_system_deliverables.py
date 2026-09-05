import os
import docx
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml import OxmlElement, parse_xml
from docx.oxml.ns import nsdecls, qn

import pptx
from pptx.util import Inches as PInches, Pt as PPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor as PRGBColor
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# 1. GENERATE FORMAL WORD DOCUMENT (.DOCX)
# ---------------------------------------------------------------------------
def set_cell_background(cell, fill_hex):
    tcPr = cell._element.get_or_add_tcPr()
    shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{fill_hex}"/>')
    tcPr.append(shd)

def set_cell_margins(cell, top=100, bottom=100, left=140, right=140):
    tcPr = cell._element.get_or_add_tcPr()
    tcMar = OxmlElement('w:tcMar')
    for m, val in [('top', top), ('bottom', bottom), ('left', left), ('right', right)]:
        node = OxmlElement(f'w:{m}')
        node.set(qn('w:w'), str(val))
        node.set(qn('w:type'), 'dxa')
        tcMar.append(node)
    tcPr.append(tcMar)

def add_styled_heading(doc, text, level):
    h = doc.add_heading(text, level=level)
    h.paragraph_format.space_before = Pt(14)
    h.paragraph_format.space_after = Pt(6)
    h.paragraph_format.keep_with_next = True
    for run in h.runs:
        run.font.name = 'Segoe UI'
        if level == 1:
            run.font.size = Pt(15)
            run.font.bold = True
            run.font.color.rgb = RGBColor(26, 54, 93) # Navy
        elif level == 2:
            run.font.size = Pt(12.5)
            run.font.bold = True
            run.font.color.rgb = RGBColor(139, 0, 0) # Crimson
        elif level == 3:
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = RGBColor(45, 55, 72)
    return h

def generate_system_word_document(output_path):
    doc = docx.Document()
    
    for section in doc.sections:
        section.top_margin = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin = Inches(1.0)
        section.right_margin = Inches(1.0)
        
    normal_style = doc.styles['Normal']
    normal_style.font.name = 'Segoe UI'
    normal_style.font.size = Pt(10.5)
    normal_style.font.color.rgb = RGBColor(45, 55, 72)
    normal_style.paragraph_format.line_spacing = 1.15
    normal_style.paragraph_format.space_after = Pt(6)

    # Header / Title Banner
    p_header = doc.add_paragraph()
    p_header.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    r_inst = p_header.add_run("UNIVERSIDAD DE SAN MARTÍN DE PORRES\nUNIDAD DE VIRTUALIZACIÓN ACADÉMICA")
    r_inst.font.size = Pt(9)
    r_inst.font.bold = True
    r_inst.font.color.rgb = RGBColor(113, 128, 150)
    
    # Main Title
    p_title = doc.add_paragraph()
    p_title.paragraph_format.space_before = Pt(10)
    p_title.paragraph_format.space_after = Pt(2)
    r_title = p_title.add_run("INFORME TÉCNICO Y PROPUESTA DE INNOVACIÓN")
    r_title.font.size = Pt(20)
    r_title.font.bold = True
    r_title.font.color.rgb = RGBColor(26, 54, 93)

    p_sub = doc.add_paragraph()
    p_sub.paragraph_format.space_after = Pt(14)
    r_sub = p_sub.add_run("Sistema Integral de Monitoreo del Cumplimiento de los Estándares de Calidad y Desempeño Docente (USMP Virtual - Pregrado y Posgrado)")
    r_sub.font.size = Pt(13)
    r_sub.font.color.rgb = RGBColor(139, 0, 0)
    r_sub.font.bold = True

    # Metadata Box
    tbl_meta = doc.add_table(rows=4, cols=2)
    tbl_meta.alignment = WD_TABLE_ALIGNMENT.CENTER
    meta_data = [
        ("Dirigido a:", "Dirección General USMP Virtual, Direcciones de Escuela, Jefaturas de Calidad y Coordinaciones"),
        ("Proyecto:", "Sistema de Monitoreo de Calidad LMS, Acompañamiento Pedagógico y BI"),
        ("Arquitectura:", "Cloud Serverless - Google Workspace (Apps Script + Sheets + Drive + Docs API)"),
        ("Alcance:", "Programas Académicos de Pregrado y Posgrado (Modalidades Virtual, Presencial e Híbrida)")
    ]
    for i, (label, val) in enumerate(meta_data):
        row = tbl_meta.rows[i]
        c0, c1 = row.cells[0], row.cells[1]
        c0.width = Inches(1.8)
        c1.width = Inches(4.7)
        set_cell_background(c0, "F0F4F8")
        set_cell_background(c1, "FFFFFF")
        set_cell_margins(c0, 60, 60, 100, 100)
        set_cell_margins(c1, 60, 60, 100, 100)
        
        p0 = c0.paragraphs[0]
        p0.paragraph_format.space_after = Pt(0)
        r0 = p0.add_run(label)
        r0.font.bold = True
        r0.font.size = Pt(9.5)
        r0.font.color.rgb = RGBColor(26, 54, 93)
        
        p1 = c1.paragraphs[0]
        p1.paragraph_format.space_after = Pt(0)
        r1 = p1.add_run(val)
        r1.font.size = Pt(9.5)
        
    doc.add_paragraph().paragraph_format.space_after = Pt(8)

    # 1. Introducción
    add_styled_heading(doc, "1. Introducción", level=1)
    doc.add_paragraph(
        "El presente documento técnico expone la justificación, arquitectura, impacto y ventajas diferenciales del "
        "Sistema de Monitoreo del Cumplimiento de los Estándares de Calidad y Desempeño Docente de la USMP Virtual. "
        "Esta plataforma integral ha sido diseñada y desarrollada para centralizar, sistematizar y automatizar el seguimiento "
        "de las aulas virtuales en las modalidades Virtual, Presencial e Híbrida, así como la supervisión del Acompañamiento "
        "Pedagógico en los programas de Pregrado y Posgrado."
    )
    doc.add_paragraph(
        "El sistema transforma un proceso históricamente fragmentado y dependiente de registros manuales en una infraestructura digital "
        "robusta, con control de concurrencia, generación desatendida de evidencias oficiales para acreditación, consolidación "
        "multidimensional de resultados y analítica de Inteligencia de Negocios (BI) en tiempo real, operando con cero costo adicional "
        "de servidores bajo el ecosistema de Google Workspace."
    )

    # 2. Requerimiento Institucional y Diagnóstico de la Problemática
    add_styled_heading(doc, "2. Requerimiento Institucional y Diagnóstico de la Problemática", level=1)
    doc.add_paragraph(
        "Los procesos de licenciamiento y acreditación universitaria (SUNEDU, SINEACE y agencias internacionales) demandan un control "
        "exhaustivo y verificable del cumplimiento de los estándares de enseñanza-aprendizaje en los entornos virtuales (LMS) y en las sesiones "
        "de tutoría y acompañamiento pedagógico. La verificación incluye rúbricas de inicio, desarrollo, cierre de asignaturas, actualización de recursos, "
        "calificación oportuna de foros/tareas, retroalimentación al estudiante y ejecución de sesiones síncronas."
    )
    doc.add_paragraph("El diagnóstico del modelo tradicional previo reveló severas ineficiencias:")

    pain_points = [
        ("Fragmentación y falta de control de concurrencia: ", "El uso de múltiples archivos de Excel locales o compartidos generaba bloqueos, sobreescrituras accidentales y desincronización entre coordinadores cuando evaluaban simultáneamente."),
        ("Altos tiempos de consolidación manual: ", "Cruzar las notas de supervisión de aulas LMS con las evaluaciones de Acompañamiento Pedagógico requería jornadas enteras de copiado, pegado y fórmulas manuales para cada semestre."),
        ("Elaboración manual de Fichas de Observación: ", "La redacción y transcripción individual de actas y fichas de cumplimiento en Word para auditorías demandaba cientos de horas hombre de digitación reiterativa."),
        ("Vulnerabilidad en la comunicación de resultados: ", "El envío manual de correos o mensajes de retroalimentación a los docentes conllevaba riesgos de adjuntar enlaces erróneos, omitir copias a coordinadores/jefaturas o utilizar formatos no oficiales."),
        ("Ceguera gerencial y falta de alertas tempranas: ", "Las direcciones y coordinaciones no disponían de un panel consolidado de KPIs para conocer en tiempo real el porcentaje de avance, el reloj de plazos (31 días) o el nivel de desempeño por escuela.")
    ]
    for p_title, p_desc in pain_points:
        bp = doc.add_paragraph(style='List Bullet')
        bp.paragraph_format.space_after = Pt(3)
        rb = bp.add_run(p_title)
        rb.font.bold = True
        rb.font.color.rgb = RGBColor(26, 54, 93)
        bp.add_run(p_desc)

    # 3. Propuesta Tecnológica Desarrollada
    add_styled_heading(doc, "3. Propuesta Tecnológica Desarrollada: Ecosistema Integral de Monitoreo", level=1)
    doc.add_paragraph(
        "La solución consiste en una plataforma web modular basada en una Single Page Application (SPA) desacoplada y Serverless, "
        "orquestada sobre Google Apps Script (GAS, motor V8), Google Sheets como base de datos transaccional con bloqueo de concurrencia "
        "(LockService), Google Drive para gestión documental y Google Docs API para la generación automatizada de actas."
    )
    doc.add_paragraph("El sistema integra 8 subsistemas especializados interconectados:")

    subsystems = [
        ("1. Seguridad, Identidad y Control de Concurrencia", "Autenticación automática por cuenta Google institucional (@usmpvirtual.edu.pe / @usmp.pe), asignación de roles (Admin, Jefe de Área, Coordinador), filtrado de filas a nivel de usuario (Row-Level Security) y cerrojo de concurrencia (LockService) para evitar colisiones entre evaluadores."),
        ("2. Gestión de Asignación y Carga Masiva", "Módulo de jefatura para asignación en bloque de cientos de asignaturas por programa académico, con cálculo visual de distribución en tiempo real con Chart.js y sincronización transversal protegida."),
        ("3. Módulos Operativos LMS (Virtual, Presencial e Híbrido)", "Indexación dinámica de rúbricas (códigos c_1_1, cp_1_1), inyección automática de marcas de tiempo de auditoría (_ts) y tolerancia a fallos en enlaces de aulas virtuales."),
        ("4. Acompañamiento Pedagógico con Semáforo de 31 Días", "Control cronológico de plazos críticos con alerta visual en 4 fases de color, cálculo vigesimal asimétrico (no castiga criterios aún no evaluados) y botones de acción mutuamente excluyentes (Felicitar vs. Reportar Fallas)."),
        ("5. Motor de Generación Automatizada de Fichas (Google Docs API)", "Clonación y llenado instantáneo de plantillas oficiales con expresiones regulares (RegEx {{c_1_1_pre}}), validación estricta de semanas completadas y guardado automático del enlace en base de datos."),
        ("6. Consolidación Multidimensional de Resultados", "DataTables de 33 columnas que cruza en memoria RAM las notas de LMS y Acompañamiento por DNI del docente, extrayendo automáticamente los criterios débiles en lenguaje natural humano y mostrándolos en pastillas de color (Pill Badges)."),
        ("7. Comunicaciones Omnicanal Automatizadas", "Emisión con un solo clic de plantillas oficiales de correo con fórmulas matemáticas vigesimales y copia automática a coordinadores y jefaturas, además de integración directa con WhatsApp Gateway (wa.me/)."),
        ("8. Inteligencia de Negocios (BI Docente y BI Coordinadores)", "Data Mart 'Sábana General' de 72 columnas, gráficos interactivos de desempeño docente por modalidad y tablero de control del equipo con clustering de tiempos reales de auditoría y productividad.")
    ]
    for s_title, s_desc in subsystems:
        p_s = doc.add_paragraph(style='List Bullet')
        p_s.paragraph_format.space_after = Pt(4)
        rs = p_s.add_run(f"{s_title}: ")
        rs.font.bold = True
        rs.font.color.rgb = RGBColor(139, 0, 0)
        p_s.add_run(s_desc)

    # 4. Justificación Cuantitativa y Tabla Comparativa
    add_styled_heading(doc, "4. Justificación Cuantitativa: Tabla Comparativa de Procesos y Tiempos", level=1)
    doc.add_paragraph(
        "A continuación se presenta la matriz de tiempos y eficiencia comparando el método manual tradicional frente a la plataforma web implementada, "
        "considerando un universo semestral de 250 asignaturas y auditorías combinadas entre Pregrado y Posgrado:"
    )

    tbl = doc.add_table(rows=7, cols=5)
    tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    headers = ["Actividad / Proceso de Monitoreo", "Método Tradicional (Manual / Excel)", "Tiempo Tradicional", "Sistema Web de Monitoreo USMP", "Tiempo con Sistema"]
    
    for c_idx, h_text in enumerate(headers):
        cell = tbl.rows[0].cells[c_idx]
        set_cell_background(cell, "1A365D")
        set_cell_margins(cell, 90, 90, 90, 90)
        p = cell.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(h_text)
        r.font.bold = True
        r.font.size = Pt(8.5)
        r.font.color.rgb = RGBColor(255, 255, 255)

    comp_rows = [
        ("1. Supervisión de Aula LMS y Registro", "Apertura de Excel local, búsqueda manual de fila, digitación de notas y cálculo de promedio.", "25 min por asignatura", "Formulario web responsivo, autocompletado, inyección de timestamp y cálculo en vivo.", "6 min por asignatura"),
        ("2. Generación de Ficha de Observación", "Descarga de plantilla Word, copia de notas una a una, formato manual y guardado en PDF.", "15 min por ficha", "Generación con 1 clic mediante Google Docs API y archivo automático en Drive.", "0.5 min (30 seg)"),
        ("3. Consolidación Multidimensional", "Cruce manual en Excel de notas LMS + Acompañamiento por DNI y ponderación vigesimal.", "12 horas por ciclo", "Cálculo instantáneo en memoria RAM y DataTables consolidado de 33 columnas.", "0 min (Instantáneo)"),
        ("4. Extracción de Criterios a Mejorar", "Revisión visual de celdas rojas, redacción manual de textos de rúbrica observados.", "8 min por docente", "Algoritmo extractor que convierte códigos a texto humano en pastillas visuales.", "0 min (Automático)"),
        ("5. Emisión de Notificaciones y Reportes", "Redacción individual de correos, adjuntos de links y envío por email o WhatsApp.", "10 min por docente", "Plantillas HTML oficiales con fórmulas vigesimales y botón directo a WhatsApp Web.", "1 min por docente"),
        ("TOTALES ESTIMADOS (Ciclo Semestral)", "Suma de horas dedicadas a tareas administrativas y de consolidación.", "~180 Horas", "Proceso automatizado y optimizado en plataforma centralizada.", "~35 Horas\n(80.5% Ahorro)")
    ]

    for r_idx, row_data in enumerate(comp_rows, start=1):
        row = tbl.rows[r_idx]
        is_total = (r_idx == len(comp_rows))
        bg = "EBF8FF" if is_total else ("F7FAFC" if r_idx % 2 == 0 else "FFFFFF")
        
        for c_idx, text in enumerate(row_data):
            cell = row.cells[c_idx]
            set_cell_background(cell, bg)
            set_cell_margins(cell, 70, 70, 70, 70)
            p = cell.paragraphs[0]
            p.paragraph_format.space_after = Pt(0)
            if c_idx in [2, 4]:
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r = p.add_run(text)
            r.font.size = Pt(8.5)
            if is_total:
                r.font.bold = True
                r.font.color.rgb = RGBColor(26, 54, 93)

    doc.add_paragraph().paragraph_format.space_after = Pt(6)

    # Callout Box
    callout_tbl = doc.add_table(rows=1, cols=1)
    callout_tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
    c_cell = callout_tbl.rows[0].cells[0]
    c_cell.width = Inches(6.5)
    set_cell_background(c_cell, "F0FFF4")
    set_cell_margins(c_cell, 120, 120, 150, 150)
    cp = c_cell.paragraphs[0]
    cp.paragraph_format.space_after = Pt(0)
    r_c1 = cp.add_run("BALANCE DE IMPACTO Y OPTIMIZACIÓN OPERATIVA:\n")
    r_c1.font.bold = True
    r_c1.font.size = Pt(10.5)
    r_c1.font.color.rgb = RGBColor(34, 84, 61)
    
    r_c2 = cp.add_run(
        "• Reducción de más de 145 horas de trabajo manual repetitivo por semestre académico.\n"
        "• Disminución del 80.5% en la carga operativa del equipo de supervisión y gestión de calidad.\n"
        "• Cero margen de error en cruces de notas, fórmulas vigesimales y asignación de cursos.\n"
        "• Cumplimiento del 100% de los requisitos documentales de auditoría en tiempo récord."
    )
    r_c2.font.size = Pt(10)
    r_c2.font.color.rgb = RGBColor(45, 55, 72)

    doc.add_paragraph().paragraph_format.space_after = Pt(6)

    # 5. Ventajas Estratégicas y Valor Agregado
    add_styled_heading(doc, "5. Ventajas Estratégicas y Diferenciales Clave", level=1)
    advs = [
        ("Gobernanza y Acreditación Continua", "Generación estandarizada de evidencias probatorias, actas y fichas de cumplimiento auditables en cualquier momento para SUNEDU, SINEACE y certificadoras internacionales."),
        ("Seguridad y Control Anti-Colisión", "Mecanismo LockService que garantiza la integridad absoluta de la base de datos aun con decenas de evaluadores concurrentes en horas pico."),
        ("Decisiones Basadas en Datos (Data-Driven)", "El Dashboard BI y la Sábana General de 72 columnas permiten a decanos y directores identificar con precisión qué competencias docentes requieren capacitaciones pedagógicas inmediatas."),
        ("Optimización del Desempeño del Equipo de Coordinadores", "El módulo de BI para coordinadores mide la dedicación efectiva, el tiempo real de auditoría por sesión (clustering inteligente) y la productividad semanal sin sesgos."),
        ("Cero Costo Adicional de Licenciamiento o Servidores", "Diseñado íntegramente sobre las cuentas de Google Workspace for Education de la USMP, eliminando costos recurrentes de infraestructura Cloud de terceros (AWS/Azure) o licencias de bases de datos propietarias.")
    ]
    for a_title, a_desc in advs:
        pa = doc.add_paragraph(style='List Bullet')
        pa.paragraph_format.space_after = Pt(4)
        ra_t = pa.add_run(f"{a_title}: ")
        ra_t.font.bold = True
        ra_t.font.color.rgb = RGBColor(139, 0, 0)
        pa.add_run(a_desc)

    # 6. Conclusiones y Recomendación
    add_styled_heading(doc, "6. Conclusiones y Recomendación Institucional", level=1)
    doc.add_paragraph(
        "1. El Sistema de Monitoreo del Cumplimiento de los Estándares de Calidad USMP - Virtual consolida un hito en la transformación digital universitaria, "
        "optimizando el 80.5% del tiempo operativo y garantizando la excelencia en la supervisión docente."
    )
    doc.add_paragraph(
        "2. La arquitectura Serverless en Google Workspace demuestra máxima eficiencia económica, alta disponibilidad, seguridad por roles y escalabilidad ilimitada "
        "para todas las escuelas de Pregrado y Posgrado."
    )
    doc.add_paragraph(
        "3. Se recomienda formalizar la plataforma como el estándar único e institucional de supervisión académica y acreditación de calidad en la modalidad virtual de la USMP."
    )

    doc.save(output_path)
    print(f"System Word document saved successfully at: {output_path}")

# ---------------------------------------------------------------------------
# 2. GENERATE POWERPOINT PRESENTATION (.PPTX)
# ---------------------------------------------------------------------------
def set_shape_flat(shape, fill_rgb, border_rgb=None, border_width=1):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if border_rgb:
        shape.line.color.rgb = border_rgb
        shape.line.width = PPt(border_width)
    else:
        shape.line.fill.background()

def create_slide_header(slide, title_text, category_text="SISTEMA DE MONITOREO DE CALIDAD - USMP VIRTUAL"):
    tb_cat = slide.shapes.add_textbox(PInches(0.8), PInches(0.4), PInches(11.5), PInches(0.35))
    tf_cat = tb_cat.text_frame
    tf_cat.word_wrap = True
    tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
    p_c = tf_cat.paragraphs[0]
    r_c = p_c.add_run()
    r_c.text = category_text.upper()
    r_c.font.size = PPt(10.5)
    r_c.font.bold = True
    r_c.font.color.rgb = PRGBColor(139, 0, 0)

    tb_title = slide.shapes.add_textbox(PInches(0.8), PInches(0.75), PInches(11.5), PInches(0.65))
    tf_title = tb_title.text_frame
    tf_title.word_wrap = True
    tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
    p_t = tf_title.paragraphs[0]
    r_t = p_t.add_run()
    r_t.text = title_text
    r_t.font.size = PPt(22)
    r_t.font.bold = True
    r_t.font.color.rgb = PRGBColor(26, 54, 93)

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, PInches(0.8), PInches(1.42), PInches(11.73), PInches(0.02)
    )
    set_shape_flat(line, PRGBColor(226, 232, 240))

def generate_system_powerpoint_presentation(output_path):
    prs = pptx.Presentation()
    prs.slide_width = PInches(13.333) # 16:9
    prs.slide_height = PInches(7.5)
    blank_layout = prs.slide_layouts[6]

    NAVY = PRGBColor(26, 54, 93)
    CRIMSON = PRGBColor(139, 0, 0)
    DARK_BLUE = PRGBColor(15, 23, 42)
    SLATE = PRGBColor(74, 85, 104)
    LIGHT_BG = PRGBColor(248, 250, 252)
    CARD_BG = PRGBColor(255, 255, 255)
    CARD_BORDER = PRGBColor(226, 232, 240)
    GREEN = PRGBColor(22, 101, 52)
    GREEN_BG = PRGBColor(240, 253, 244)

    card_w = PInches(2.7)
    card_h = PInches(4.8)
    spacing = PInches(0.3)
    start_x = PInches(0.8)
    y_pos = PInches(1.8)

    # -------------------------------------------------------------
    # SLIDE 1: PORTADA
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(13.333), PInches(7.5))
    set_shape_flat(bg1, NAVY)

    accent_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(0.4), PInches(7.5))
    set_shape_flat(accent_bar, CRIMSON)

    tb_title = s1.shapes.add_textbox(PInches(1.2), PInches(1.4), PInches(10.8), PInches(4.8))
    tf1 = tb_title.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "UNIVERSIDAD DE SAN MARTÍN DE PORRES  |  USMP VIRTUAL\n"
    r0.font.size = PPt(13.5)
    r0.font.bold = True
    r0.font.color.rgb = PRGBColor(226, 232, 240)

    p1 = tf1.add_paragraph()
    p1.space_before = PPt(14)
    r1 = p1.add_run()
    r1.text = "Sistema de Monitoreo del Cumplimiento de los Estándares de Calidad y Desempeño Docente"
    r1.font.size = PPt(30)
    r1.font.bold = True
    r1.font.color.rgb = PRGBColor(255, 255, 255)

    p2 = tf1.add_paragraph()
    p2.space_before = PPt(12)
    r2 = p2.add_run()
    r2.text = "Automatización Integral, Acreditación Continua y Business Intelligence (Pregrado y Posgrado)"
    r2.font.size = PPt(16.5)
    r2.font.color.rgb = PRGBColor(214, 158, 46)

    p3 = tf1.add_paragraph()
    p3.space_before = PPt(28)
    r3 = p3.add_run()
    r3.text = "Arquitectura Serverless: Google Apps Script  •  Google Sheets  •  Google Drive  •  Docs API\nUnidad de Virtualización Académica  •  2026"
    r3.font.size = PPt(12.5)
    r3.font.color.rgb = PRGBColor(203, 213, 225)

    # -------------------------------------------------------------
    # SLIDE 2: PROBLEMÁTICA TRADICIONAL
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    bg2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(13.333), PInches(7.5))
    set_shape_flat(bg2, LIGHT_BG)
    create_slide_header(s2, "Diagnóstico: Cuellos de Botella del Monitoreo Tradicional", "PROBLEMATICA INICIAL")

    problems = [
        ("Sobreescritura y Bloqueos", "Uso de Excels compartidos sin control de concurrencia. Pérdida de evaluaciones y colisión entre coordinadores."),
        ("Cruce Manual Complejo", "Jornadas enteras cruzando manualmente notas LMS y Acompañamiento por DNI con alto riesgo de error humano."),
        ("Fichas Word Manuales", "Redacción artesanal de actas y fichas de acreditación (~15 min por curso), sumando cientos de horas hombre."),
        ("Falta de BI en Tiempo Real", "Desconocimiento de indicadores clave de avance y debilidades pedagógicas hasta el cierre de periodo.")
    ]

    for i, (p_title, p_desc) in enumerate(problems):
        x = start_x + i * (card_w + spacing)
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, card_w, card_h)
        set_shape_flat(card, CARD_BG, CARD_BORDER)

        pill = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + PInches(0.2), y_pos + PInches(0.25), PInches(0.5), PInches(0.5))
        set_shape_flat(pill, CRIMSON)
        tf_p = pill.text_frame
        p_num = tf_p.paragraphs[0]
        r_num = p_num.add_run()
        r_num.text = f"0{i+1}"
        r_num.font.bold = True
        r_num.font.size = PPt(14)
        r_num.font.color.rgb = PRGBColor(255, 255, 255)
        p_num.alignment = PP_ALIGN.CENTER

        tb = s2.shapes.add_textbox(x + PInches(0.2), y_pos + PInches(0.9), card_w - PInches(0.4), card_h - PInches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        pt = tf.paragraphs[0]
        rt = pt.add_run()
        rt.text = p_title
        rt.font.size = PPt(14.5)
        rt.font.bold = True
        rt.font.color.rgb = NAVY

        pd = tf.add_paragraph()
        pd.space_before = PPt(10)
        rd = pd.add_run()
        rd.text = p_desc
        rd.font.size = PPt(12)
        rd.font.color.rgb = SLATE

    # -------------------------------------------------------------
    # SLIDE 3: ARQUITECTURA DE SUBSISTEMAS
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    bg3 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(13.333), PInches(7.5))
    set_shape_flat(bg3, LIGHT_BG)
    create_slide_header(s3, "La Solución: Ecosistema Modular Integrado (8 Subsistemas)", "ARQUITECTURA DE SOFTWARE")

    sub_cards = [
        ("Seguridad y LockService", "Control de concurrencia y roles estrictos. Filtrado por fila (Row-Level Security) según coordinador."),
        ("Operativo LMS y Semáforo", "Indexación dinámica de rúbricas (Virtual/Presencial) y reloj de 31 días en Acompañamiento Pedagógico."),
        ("Motor de Fichas (Docs API)", "Generación instantánea de actas oficiales con RegEx y almacenamiento desatendido en Google Drive."),
        ("Business Intelligence", "Data Mart 'Sábana General' de 72 cols, gráficos dinámicos y BI de productividad de coordinadores.")
    ]

    for i, (a_title, a_desc) in enumerate(sub_cards):
        x = start_x + i * (card_w + spacing)
        card = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, card_w, card_h)
        set_shape_flat(card, CARD_BG, CARD_BORDER)

        pill = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + PInches(0.2), y_pos + PInches(0.25), PInches(0.5), PInches(0.5))
        set_shape_flat(pill, NAVY)
        tf_p = pill.text_frame
        p_num = tf_p.paragraphs[0]
        r_num = p_num.add_run()
        r_num.text = "✓"
        r_num.font.bold = True
        r_num.font.size = PPt(16)
        r_num.font.color.rgb = PRGBColor(255, 255, 255)
        p_num.alignment = PP_ALIGN.CENTER

        tb = s3.shapes.add_textbox(x + PInches(0.2), y_pos + PInches(0.9), card_w - PInches(0.4), card_h - PInches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        pt = tf.paragraphs[0]
        rt = pt.add_run()
        rt.text = a_title
        rt.font.size = PPt(14.5)
        rt.font.bold = True
        rt.font.color.rgb = NAVY

        pd = tf.add_paragraph()
        pd.space_before = PPt(10)
        rd = pd.add_run()
        rd.text = a_desc
        rd.font.size = PPt(12)
        rd.font.color.rgb = SLATE

    # -------------------------------------------------------------
    # SLIDE 4: CUADRO COMPARATIVO DE TIEMPOS
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    bg4 = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(13.333), PInches(7.5))
    set_shape_flat(bg4, LIGHT_BG)
    create_slide_header(s4, "Análisis Comparativo de Procesos y Tiempos de Supervisión", "EFICIENCIA OPERATIVA")

    t_shape = s4.shapes.add_table(6, 5, PInches(0.8), PInches(1.8), PInches(11.73), PInches(4.8))
    table = t_shape.table
    table.columns[0].width = PInches(2.5)
    table.columns[1].width = PInches(3.1)
    table.columns[2].width = PInches(1.4)
    table.columns[3].width = PInches(3.3)
    table.columns[4].width = PInches(1.43)

    t_headers = ["Proceso de Monitoreo", "Método Tradicional (Excel/Word)", "Tiempo Trad.", "Sistema Web USMP", "Tiempo Sistema"]
    for c_idx, h in enumerate(t_headers):
        cell = table.cell(0, c_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = h
        r.font.bold = True
        r.font.size = PPt(11.5)
        r.font.color.rgb = PRGBColor(255, 255, 255)

    rows_data = [
        ("1. Supervisión de Aula LMS", "Búsqueda en Excel, digitación y fórmulas manuales.", "25 min", "Formulario web responsivo, autocompletado y timestamps.", "6 min"),
        ("2. Ficha de Observación", "Copia manual de notas a plantilla Word y exportación a PDF.", "15 min", "Generación con 1 clic con Google Docs API a Drive.", "30 seg (Auto)"),
        ("3. Consolidación de Notas", "Cruce manual de notas LMS + Acompañamiento por DNI.", "12 horas", "Cálculo en memoria RAM y DataTables de 33 columnas.", "0 min (Auto)"),
        ("4. Extracción de Fallas", "Revisión celda por celda y redacción manual de textos.", "8 min / doc.", "Extractor automático que convierte códigos en texto humano.", "0 min (Auto)"),
        ("5. Notificación y Reporte", "Redacción de correos individuales y copia manual.", "10 min / doc.", "Plantilla oficial con fórmulas vigesimales y WhatsApp Web.", "1 min")
    ]

    for r_idx, r_data in enumerate(rows_data, start=1):
        bg_c = CARD_BG if r_idx % 2 != 0 else PRGBColor(241, 245, 249)
        for c_idx, val in enumerate(r_data):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_c
            p = cell.text_frame.paragraphs[0]
            if c_idx in [2, 4]:
                p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = val
            r.font.size = PPt(11)
            if c_idx == 4:
                r.font.bold = True
                r.font.color.rgb = GREEN
            elif c_idx == 2:
                r.font.bold = True
                r.font.color.rgb = CRIMSON
            else:
                r.font.color.rgb = DARK_BLUE

    # -------------------------------------------------------------
    # SLIDE 5: BIG STATS DE IMPACTO
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    bg5 = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(13.333), PInches(7.5))
    set_shape_flat(bg5, LIGHT_BG)
    create_slide_header(s5, "Impacto Cuantitativo: 80.5% de Reducción en Carga Operativa", "MÉTRICAS CLAVE")

    m_card_w = PInches(3.64)
    m_card_h = PInches(4.8)
    m_spacing = PInches(0.4)
    m_start_x = PInches(0.8)

    metrics = [
        ("-80.5%", "Optimización de Tiempo Total", "De ~180 horas a solo 35 horas totales por semestre en tareas de supervisión y gestión.", GREEN, GREEN_BG),
        ("145+ Horas", "Horas Hombre Recuperadas", "Tiempo reasignado a acompañamiento pedagógico proactivo y mejora continua del docente.", NAVY, PRGBColor(235, 248, 255)),
        ("1 Clic", "Generación de Fichas Oficiales", "Creación de actas con formato institucional listas para auditorías de licenciamiento y acreditación.", CRIMSON, PRGBColor(255, 245, 245))
    ]

    for i, (stat, subtitle, desc, col, bg_col) in enumerate(metrics):
        x = m_start_x + i * (m_card_w + m_spacing)
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, m_card_w, m_card_h)
        set_shape_flat(card, bg_col, col, 2)

        tb = s5.shapes.add_textbox(x + PInches(0.3), y_pos + PInches(0.5), m_card_w - PInches(0.6), m_card_h - PInches(1.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        p_s = tf.paragraphs[0]
        p_s.alignment = PP_ALIGN.CENTER
        r_s = p_s.add_run()
        r_s.text = stat
        r_s.font.size = PPt(42)
        r_s.font.bold = True
        r_s.font.color.rgb = col

        p_sub = tf.add_paragraph()
        p_sub.alignment = PP_ALIGN.CENTER
        p_sub.space_before = PPt(10)
        r_sub = p_sub.add_run()
        r_sub.text = subtitle
        r_sub.font.size = PPt(15.5)
        r_sub.font.bold = True
        r_sub.font.color.rgb = DARK_BLUE

        p_d = tf.add_paragraph()
        p_d.alignment = PP_ALIGN.CENTER
        p_d.space_before = PPt(14)
        r_d = p_d.add_run()
        r_d.text = desc
        r_d.font.size = PPt(12)
        r_d.font.color.rgb = SLATE

    # -------------------------------------------------------------
    # SLIDE 6: BUSINESS INTELLIGENCE
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    bg6 = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(13.333), PInches(7.5))
    set_shape_flat(bg6, LIGHT_BG)
    create_slide_header(s6, "Business Intelligence: Decisiones Académicas en Tiempo Real", "INTELIGENCIA DE NEGOCIOS")

    bi_cards = [
        ("Sábana General (72 Cols)", "Data Mart estructurado que unifica asignación, LMS expandido, acompañamiento y score vigesimal final."),
        ("Doble Gráfico Dinámico", "Distribución de niveles (Doughnut) y promedios por criterio (Barras) adaptados a la modalidad elegida."),
        ("Productividad del Equipo", "Algoritmo de clustering que calcula el tiempo neto real de auditoría LMS descartando tiempos muertos."),
        ("Pastillas Visuales (Pills)", "Identificación inmediata en el consolidado de los criterios pedagógicos que requieren refuerzo.")
    ]

    for i, (title, desc) in enumerate(bi_cards):
        x = start_x + i * (card_w + spacing)
        card = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, card_w, card_h)
        set_shape_flat(card, CARD_BG, CARD_BORDER)

        pill = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + PInches(0.2), y_pos + PInches(0.25), PInches(0.5), PInches(0.5))
        set_shape_flat(pill, NAVY)
        tf_p = pill.text_frame
        p_num = tf_p.paragraphs[0]
        r_num = p_num.add_run()
        r_num.text = "📈"
        r_num.font.bold = True
        r_num.font.size = PPt(15)
        p_num.alignment = PP_ALIGN.CENTER

        tb = s6.shapes.add_textbox(x + PInches(0.2), y_pos + PInches(0.9), card_w - PInches(0.4), card_h - PInches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        pt = tf.paragraphs[0]
        rt = pt.add_run()
        rt.text = title
        rt.font.size = PPt(14.5)
        rt.font.bold = True
        rt.font.color.rgb = NAVY

        pd = tf.add_paragraph()
        pd.space_before = PPt(10)
        rd = pd.add_run()
        rd.text = desc
        rd.font.size = PPt(12)
        rd.font.color.rgb = SLATE

    # -------------------------------------------------------------
    # SLIDE 7: CONCLUSIONES
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    bg7 = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PInches(13.333), PInches(7.5))
    set_shape_flat(bg7, NAVY)

    tb_c = s7.shapes.add_textbox(PInches(1.2), PInches(1.0), PInches(10.8), PInches(5.5))
    tfc = tb_c.text_frame
    tfc.word_wrap = True

    pc_title = tfc.paragraphs[0]
    rc_title = pc_title.add_run()
    rc_title.text = "Conclusiones y Recomendación Institucional"
    rc_title.font.size = PPt(28)
    rc_title.font.bold = True
    rc_title.font.color.rgb = PRGBColor(255, 255, 255)

    conclusions = [
        "1. Innovación Tecnológica de Alto Impacto: Reemplaza hojas de cálculo dispersas por una arquitectura Cloud Serverless segura y escalable.",
        "2. Retorno Inmediato en Eficiencia (80.5%): Ahorro directo de más de 145 horas hombre por semestre, liberando al equipo para labores de calidad docente.",
        "3. Solidez en Acreditación y Licenciamiento: Evidencias documentales oficiales generadas al instante con trazabilidad inmutable.",
        "4. Cero Costo Adicional: 100% operativo sobre las licencias de Google Workspace institucionales de la USMP.",
        "RECOMENDACIÓN: Institucionalizar la plataforma como el sistema oficial obligatorio de supervisión académica y acreditación en Pregrado y Posgrado."
    ]

    for idx, c_text in enumerate(conclusions):
        p_item = tfc.add_paragraph()
        p_item.space_before = PPt(14 if idx < 4 else 20)
        r_item = p_item.add_run()
        r_item.text = c_text
        if idx == 4:
            r_item.font.size = PPt(15)
            r_item.font.bold = True
            r_item.font.color.rgb = PRGBColor(214, 158, 46)
        else:
            r_item.font.size = PPt(13)
            r_item.font.color.rgb = PRGBColor(226, 232, 240)

    prs.save(output_path)
    print(f"System PowerPoint presentation saved successfully at: {output_path}")

if __name__ == "__main__":
    base_dir = r"c:\Proyectos  Antigravity\Sietama de monitoreo pregrado"
    doc_path = os.path.join(base_dir, "Informe_Tecnico_Sistema_Monitoreo_Calidad_USMP.docx")
    ppt_path = os.path.join(base_dir, "Presentacion_Sistema_Monitoreo_Calidad_USMP.pptx")
    
    generate_system_word_document(doc_path)
    generate_system_powerpoint_presentation(ppt_path)
